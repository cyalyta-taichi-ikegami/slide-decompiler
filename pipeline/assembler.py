"""python-pptx によるスライド再構成（PPTX 組み立て）レイヤー

修復済み背景画像をスライド全面に敷き、
その上に翻訳済みテキストボックスを正確な EMU 座標で重ねる。
"""
import io
import logging

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from config import SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU
from models.slide_elements import SlideAnalysis
from utils.coordinate_mapper import normalized_to_emu, estimate_font_size_pt

logger = logging.getLogger(__name__)


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """#RRGGBB → (R, G, B)"""
    h = hex_color.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def build_slide(
    prs: Presentation,
    background_image_bytes: bytes,
    analysis: SlideAnalysis,
) -> None:
    """既存の Presentation に 1 スライドを追加する。

    レイヤー構造:
        Layer 1（最背面）: background_image_bytes をスライド全面に配置
        Layer 2（前面）  : analysis の is_editable=True 要素をテキストボックスで配置

    Args:
        prs:                      追加先の Presentation オブジェクト（破壊的変更あり）
        background_image_bytes:   インペインティング済みの背景画像（PNG bytes）
        analysis:                 翻訳済みの SlideAnalysis
    """
    slide_layout = prs.slide_layouts[6]  # ブランクレイアウト
    slide = prs.slides.add_slide(slide_layout)

    # --- Layer 1: 背景画像をスライド全面に配置 ---
    img_stream = io.BytesIO(background_image_bytes)
    slide.shapes.add_picture(
        img_stream,
        left=Emu(0),
        top=Emu(0),
        width=Emu(SLIDE_WIDTH_EMU),
        height=Emu(SLIDE_HEIGHT_EMU),
    )

    # --- Layer 2: テキストボックスを各座標に配置 ---
    editable = [e for e in analysis.elements if e.is_editable and e.text_content.strip()]

    for element in editable:
        left, top, width, height = normalized_to_emu(element.bounding_box)
        font_pt = estimate_font_size_pt(element.bounding_box)
        r, g, b = _hex_to_rgb(element.color_hex)

        txBox = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = element.text_content

        font = run.font
        font.size = Pt(font_pt)
        font.color.rgb = RGBColor(r, g, b)
        font.bold = element.font_bold

        # 翻訳で文字列が長くなった場合もボックス内に収まるよう自動縮小する
        try:
            tf.fit_text(max_size=font_pt, min_size=6)
        except Exception:
            pass  # fit_text が失敗しても処理は続行する

    logger.info(f"{len(editable)} 個のテキストボックスを配置")
